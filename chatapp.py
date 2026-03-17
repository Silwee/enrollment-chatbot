import streamlit as st
import uuid
import io
from google import genai
from google.genai import types
from google.genai.types import Part
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from openpyxl import load_workbook

API_KEY = "AIzaSyB2pD4hHe1kJp20Mtw3468U9CQAb3X8mts"

SUMMARIZE_PROMPT = (
    "Tóm tắt đoạn văn bản hoặc file sau dưới dạng danh sách các nội dung quan trọng của văn bản. "
    "Không cần phải giữ lại tất cả các chi tiết, nhưng đừng bỏ qua bất kỳ thông tin quan trọng nào. "
    "Hãy tóm tắt một cách ngắn gọn và dễ hiểu, sử dụng các gạch đầu dòng nếu cần."
    "Nếu có thể, hãy cố gắng giữ lại các con số, tên riêng, và các chi tiết cụ thể khác mà có thể hữu ích cho việc trả lời câu hỏi sau này."
)

SYSTEM_INSTRUCTION = (
    "Bạn là 1 nhân viên hỗ trợ tuyển sinh cho 1 trường học. "
    "Dung lượng kiến thức của bạn là tất cả các tài liệu được đính kèm dưới đây. "
    "Hãy trả lời ngắn gọn và chính xác dựa trên các tài liệu đó, "
    "không suy đoán nếu tài liệu không cung cấp thông tin."
)


# ---------------------------------------------------------------------------
# Shared in-memory store — one instance across ALL Streamlit sessions
# ---------------------------------------------------------------------------
@st.cache_resource
def get_summaries_store() -> dict:
    """Returns the shared dict: { filename: summary_text }"""
    return {}


# ---------------------------------------------------------------------------
# Local text extraction for formats unsupported by the Files API
# ---------------------------------------------------------------------------
def extract_text(file_name: str, file_bytes: bytes, mime_type: str) -> str | None:
    """Return plain text from the file, or None if the file should be uploaded directly."""
    ext = file_name.rsplit(".", 1)[-1].lower()

    if ext in ("docx", "doc"):
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if ext in ("xlsx", "xls"):
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                line = "\t".join(str(c) for c in row if c is not None)
                if line.strip():
                    lines.append(line)
        return "\n".join(lines)

    if ext in ("pptx", "ppt"):
        prs = PptxPresentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            lines.append(text)
        return "\n".join(lines)

    if ext == "txt" or mime_type.startswith("text/"):
        return file_bytes.decode("utf-8", errors="replace")

    # PDF and other Files-API-supported types → upload directly
    return None


# ---------------------------------------------------------------------------
# Summarize a file using the Gemini Files API
# ---------------------------------------------------------------------------
def summarize_file(client: genai.Client, file_name: str, file_bytes: bytes, mime_type: str) -> str:
    """Extract or upload file content, ask the model to summarize it, return summary text."""
    text = extract_text(file_name, file_bytes, mime_type)

    if text is not None:
        # Send extracted plain text directly — no upload needed
        response = client.models.generate_content(
            model="gemini-2.5-flash-lite",
            contents=f"{SUMMARIZE_PROMPT}\n\n{text}",
        )
    else:
        # Upload via Files API (e.g. PDF)
        uploaded = client.files.upload(
            file=io.BytesIO(file_bytes),
            config={"display_name": file_name, "mime_type": mime_type},
        )
        response = None
        try:
            response = client.models.generate_content(
                model="gemini-2.5-flash-lite",
                contents=[
                    Part.from_uri(file_uri=uploaded.uri, mime_type=uploaded.mime_type),
                    Part(text=SUMMARIZE_PROMPT),
                ],
            )
        finally:
            try:
                client.files.delete(name=uploaded.name)
            except Exception:
                pass
        if response is None:
            raise RuntimeError("No response received from the model.")

    return response.text


# ---------------------------------------------------------------------------
# Main Application
# ---------------------------------------------------------------------------
def main():
    st.set_page_config(page_title="Chatbot AI hỗ trợ tuyển sinh", page_icon="🎓", layout="wide")

    # Shared summaries store (persists for the lifetime of the server process)
    summaries: dict = get_summaries_store()

    # Gemini client
    try:
        client = genai.Client(api_key=API_KEY)
    except Exception:
        st.error("Could not initialise Gemini client. Check the API key.")
        st.stop()

    # Per-tab session state
    if "session_token" not in st.session_state:
        st.session_state.session_token = str(uuid.uuid4())
        st.session_state.messages = []
    if "admin_authenticated" not in st.session_state:
        st.session_state.admin_authenticated = False

    # -----------------------------------------------------------------------
    # Layout: sidebar for document management, main area for chat
    # -----------------------------------------------------------------------
    with st.sidebar:
        st.header("📂 Quản lý tài liệu")
        st.caption(f"Session: `{st.session_state.session_token[:8]}…`")

        # --- Auth gate ---
        if not st.session_state.admin_authenticated:
            st.divider()
            st.subheader("🔒 Đăng nhập quản trị")
            with st.form("admin_login_form"):
                username = st.text_input("Tên đăng nhập")
                password = st.text_input("Mật khẩu", type="password")
                submitted = st.form_submit_button("Đăng nhập", use_container_width=True)
                if submitted:
                    if username == "admin" and password == "admin":
                        st.session_state.admin_authenticated = True
                        st.rerun()
                    else:
                        st.error("Tên đăng nhập hoặc mật khẩu không đúng.")
        else:
            # Logged-in: show logout button
            if st.button("Đăng xuất", use_container_width=True):
                st.session_state.admin_authenticated = False
                st.rerun()

            # --- File uploader ---
            uploaded_files = st.file_uploader(
                "Tải lên tài liệu (PDF, Word, Excel, PowerPoint, TXT)…",
                accept_multiple_files=True,
                type=["pdf", "docx", "doc", "txt", "xlsx", "pptx"],
                help="Tải lên file để làm nội dung cho chatbot AI.",
            )

            if uploaded_files:
                new_files = [f for f in uploaded_files if f.name not in summaries]
                if new_files:
                    for uf in new_files:
                        with st.spinner(f"Đang tóm tắt **{uf.name}**…"):
                            mime = uf.type or "application/octet-stream"
                            try:
                                summary = summarize_file(client, uf.name, uf.read(), mime)
                                summaries[uf.name] = summary
                                st.success(f"✅ {uf.name}")
                            except Exception as e:
                                st.error(f"❌ {uf.name}: {e}")
                else:
                    st.info("Tất cả các file đã được tóm tắt.")

            # --- Display all summaries ---
            st.divider()
            st.subheader(f"📋 Nội dung ({len(summaries)} văn bản)",
                         help="Nội dung tóm tắt được sử dụng làm kiến thức nền cho chatbot AI. Nhấn vào từng văn bản để xem nội dung, chỉnh sửa hoặc xoá.")

            if not summaries:
                st.caption("Chưa có tài liệu nào. Vui lòng tải tài liệu ở trên.")
            else:
                for fname in list(summaries.keys()):
                    with st.expander(f"📄 {fname}"):
                        edited = st.text_area(
                            label="Chỉnh sửa",
                            value=summaries[fname],
                            height=300,
                            key=f"summary_edit_{fname}",
                            label_visibility="collapsed",
                        )
                        col1, col2 = st.columns([1, 1])
                        is_edited = False
                        with col1:
                            if st.button("💾 Lưu", key=f"save_{fname}"):
                                summaries[fname] = edited
                                is_edited = True
                        with col2:
                            if st.button("🗑️ Xoá", key=f"remove_{fname}"):
                                del summaries[fname]
                                st.rerun()
                        if is_edited:
                            st.success("Lưu thành công!")

    # -----------------------------------------------------------------------
    # Main chat area
    # -----------------------------------------------------------------------
    st.title("🎓 Chatbot hỗ trợ tuyển sinh")

    # Render chat history
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # New user input
    if prompt := st.chat_input("VD: Thi lớp 10 vào ngày nào?"):

        st.chat_message("user").markdown(prompt)
        st.session_state.messages.append({"role": "user", "content": prompt})

        # Build the full context: system instruction + all summaries
        if summaries:
            docs_context = "\n\n".join(
                f"### {fname}\n{summary}" for fname, summary in summaries.items()
            )
            full_system = f"{SYSTEM_INSTRUCTION}\n\n---\n\n## Tài liệu tham khảo:\n\n{docs_context}"
        else:
            full_system = SYSTEM_INSTRUCTION

        # Build conversation history for the API
        contents = []
        for msg in st.session_state.messages:
            contents.append(
                types.Content(
                    role="user" if msg["role"] == "user" else "model",
                    parts=[Part(text=msg["content"])],
                )
            )

        with st.chat_message("model"):
            with st.spinner("Đang suy nghĩ…"):
                response = client.models.generate_content(
                    model="gemini-2.5-flash-lite",
                    contents=contents,
                    config=types.GenerateContentConfig(
                        system_instruction=full_system,
                        temperature=0.3,
                    ),
                )
                st.markdown(response.text)

        st.session_state.messages.append({"role": "model", "content": response.text})


if __name__ == "__main__":
    main()



